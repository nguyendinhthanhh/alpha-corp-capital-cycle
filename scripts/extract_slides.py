import os
import glob
from pptx import Presentation

materials_dir = os.path.join(os.path.dirname(__file__), '../docs/course-materials/original')
output_dir = os.path.join(os.path.dirname(__file__), '../docs/quiz')
extracted_dir = os.path.join(output_dir, 'extracted_text')

if not os.path.exists(output_dir):
    os.makedirs(output_dir)
if not os.path.exists(extracted_dir):
    os.makedirs(extracted_dir)

files = glob.glob(os.path.join(materials_dir, '*.*'))

manifest = []

for filepath in files:
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1].lower()
    
    # Construct Source ID
    parts = filename.split(' ')
    prefix = parts[0]
    num = parts[1] if len(parts) > 1 else ''
    
    prefix = prefix.replace('Session', 'S').replace('Slot', 'SL')
    source_id = f"MLN122-{prefix}{num}"

    # Determine Chapter and Session
    import re
    chapter_match = re.search(r'Ch[uưươ]+ng\s*(\d)', filename, re.IGNORECASE)
    chapter = chapter_match.group(1) if chapter_match else 'Unknown'
    
    session_match = re.search(r'(Session|Slot)\s*(\d+)', filename, re.IGNORECASE)
    session = f"{session_match.group(1)} {session_match.group(2)}" if session_match else 'Unknown'

    status = 'readable'
    note = ''
    slide_count = 0
    text_content = ''

    if ext == '.ppt':
        status = 'UNREADABLE_SOURCE'
        note = 'Needs conversion to .pptx or .pdf (old .ppt format is not supported by parser).'
        slide_count = 'N/A'
    elif ext == '.pptx':
        try:
            prs = Presentation(filepath)
            slide_count = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                text_content += f"\n--- Slide {i + 1} ---\n"
                
                # Try to get notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        text_content += f"[Speaker Notes]\n{notes}\n\n"
                        
                # Get slide text
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += f"{shape.text}\n"
                        
        except Exception as e:
            status = 'UNREADABLE_SOURCE'
            note = f"Parsing error: {str(e)}"
            slide_count = 'N/A'
    else:
        status = 'UNREADABLE_SOURCE'
        note = f"Unsupported extension: {ext}"
        slide_count = 'N/A'

    if status == 'readable':
        out_path = os.path.join(extracted_dir, f"{filename}.txt")
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(text_content)

    manifest.append({
        'SourceID': source_id,
        'File': filename,
        'Chapter': chapter,
        'Session': session,
        'Status': status,
        'Slides': slide_count,
        'Topics': 'To be extracted',
        'Note': note
    })

# Write Manifest
manifest_path = os.path.join(output_dir, 'SOURCE_MANIFEST.md')
with open(manifest_path, 'w', encoding='utf-8') as f:
    f.write('# MLN122 Source Manifest\n\n')
    f.write('| Source ID | File | Chapter | Session | Status | Slides | Note |\n')
    f.write('|---|---|---|---|---|---|---|\n')
    for item in manifest:
        f.write(f"| {item['SourceID']} | {item['File']} | {item['Chapter']} | {item['Session']} | {item['Status']} | {item['Slides']} | {item['Note']} |\n")

print(f"Extraction complete. Manifest saved to {manifest_path}")
